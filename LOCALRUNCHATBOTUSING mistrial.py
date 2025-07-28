# Run with: python -m streamlit run "C:/Users/soumyadip.mukherjee/Desktop/Optimized_Chatbot.py"
# Prerequisites: pip install streamlit pymupdf pytesseract pillow python-docx python-pptx faiss-cpu langchain

import streamlit as st
import os, glob, re, shutil, random, smtplib
from email.message import EmailMessage
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.llms import Ollama
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain_core.documents import Document

st.set_page_config(page_title="🛡️ CLG Private Chatbot", layout="wide")

for key in ["page", "is_authenticated", "generated_otp", "email", "history"]:
    if key not in st.session_state:
        st.session_state[key] = {
            "page": "login",
            "is_authenticated": False,
            "generated_otp": None,
            "email": "",
            "history": []
        }[key]

# -------------------- Utility Functions --------------------
def send_otp(email):
    otp = str(random.randint(100000, 999999))
    st.session_state.generated_otp = otp

    sender = "mukherjeesam8129@gmail.com"
    app_password = "pvmf msqz apow ddii"  # WARNING: Replace with secure method

    msg = EmailMessage()
    msg["Subject"] = "Your OTP for CLG Chatbot"
    msg["From"] = sender
    msg["To"] = email
    msg.set_content(f"Your OTP is: {otp}")

    try:
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
            smtp.login(sender, app_password)
            smtp.send_message(msg)
        st.success(f"✅ OTP sent to {email}")
    except Exception as e:
        st.error(f"❌ Failed to send OTP: {e}")

def sanitize_filenames(folder):
    for filename in os.listdir(folder):
        clean_name = re.sub(r'[\\/*?:"<>|\[\]]+', "_", filename)
        shutil.move(os.path.join(folder, filename), os.path.join(folder, clean_name))

def load_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        extracted = page.get_text()
        if not extracted.strip():
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            extracted = pytesseract.image_to_string(img)
        text += extracted + "\n"
    return [Document(page_content=text)]

def load_docx(file_path):
    doc = DocxDocument(file_path)
    return [Document(page_content="\n".join(p.text for p in doc.paragraphs))]

def load_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return [Document(page_content=text)]
    except PackageNotFoundError:
        st.warning(f"⚠️ Skipping unreadable PPTX: {os.path.basename(file_path)}")
        return []

def load_image(file_path):
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return [Document(page_content=text)]
    except Exception as e:
        st.warning(f"⚠️ Could not process image {file_path}: {e}")
        return []

def load_files(folder):
    sanitize_filenames(folder)
    docs = []
    for file_path in glob.glob(os.path.join(folder, "*.*")):
        ext = file_path.lower().split(".")[-1]
        if ext == "pdf":
            docs += load_pdf(file_path)
        elif ext == "docx":
            docs += load_docx(file_path)
        elif ext == "pptx":
            docs += load_pptx(file_path)
        elif ext in ["jpg", "jpeg", "png"]:
            docs += load_image(file_path)
    return docs

# -------------------- Login Page --------------------
def login_page():
    st.title("🔐 Login to CLG Chatbot")
    username = st.text_input("Enter your username", placeholder="E.g: soumyadip.mukherjee")
    if username:
        st.session_state.email = username + "@vit.ac.in"

    if st.button("Send OTP"):
        if username:
            send_otp(st.session_state.email)
        else:
            st.error("Username required")

    if st.session_state.generated_otp:
        otp = st.text_input("Enter the OTP", type="password")
        if st.button("Verify OTP"):
            if otp == st.session_state.generated_otp:
                st.session_state.is_authenticated = True
                st.session_state.page = "chatbot"
                st.success("✅ Login successful")
            else:
                st.error("❌ Incorrect OTP")

# -------------------- Chatbot Page --------------------
def chatbot_page():
    st.title("🤖 CLG Private ChatBot")

    if st.button("Logout"):
        st.session_state.clear()
        st.session_state.page = "login"
        return

    folder_path = "C:/Users/soumyadip.mukherjee/Downloads/hcp"
    os.makedirs(folder_path, exist_ok=True)

    st.sidebar.header("📁 Load Data from Folder")
    uploaded_files = st.sidebar.file_uploader("Attach file(s)", type=["pdf", "docx", "pptx", "jpg", "jpeg", "png"], accept_multiple_files=True)

    for uploaded_file in uploaded_files:
        with open(os.path.join(folder_path, uploaded_file.name), "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.sidebar.success(f"Uploaded: {uploaded_file.name}")

    docs = load_files(folder_path)

    if docs:
        with st.spinner("🔍 Processing documents..."):
            splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
            chunks = splitter.split_documents(docs)

            try:
                embeddings = OllamaEmbeddings(model="nomic-embed-text", base_url="http://localhost:11434")
                vectordb = FAISS.from_documents(chunks, embeddings)
                retriever = vectordb.as_retriever()
            except Exception as e:
                st.error(f"Vector DB error: {e}")
                return

        simple_keywords = ["what is", "define", "who is", "list", "name"]
        complex_keywords = ["explain", "compare", "difference", "how does", "summarize", "step by step", "generate code"]

        def choose_model(user_input):
            input_lower = user_input.lower()
            if any(kw in input_lower for kw in complex_keywords):
                return "llama3:instruct"
            return "mistral"

        st.subheader("💬 Ask Your Questions")
        with st.form("chat_form", clear_on_submit=True):
            user_input = st.text_input("Type your question:")
            submitted = st.form_submit_button("Send")

        if submitted and user_input:
            with st.spinner("🤖 Thinking..."):
                try:
                    chosen_model = choose_model(user_input)
                    st.info(f"Using model: `{chosen_model}`")

                    llm = Ollama(model=chosen_model, base_url="http://localhost:11434")
                    qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)

                    answer = qa_chain.run(user_input)
                    if not answer or "i don't know" in answer.lower():
                        answer += "\n\n🔗 [Contact support](https://www.vit.ac.in.com/support)"
                    st.session_state.history.append((user_input, answer))
                except Exception as e:
                    st.session_state.history.append((user_input, f"❗ Error: {e}"))

        st.markdown('<div class="chat-container">', unsafe_allow_html=True)
        for user_msg, bot_msg in reversed(st.session_state.history):
            st.markdown(f'<div class="message user">🧑 You: {user_msg}</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="message bot">🤖 Bot: {bot_msg}</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

        st.markdown("""
        <style>
            .message.user {
                background-color: #f0f0f0;
                color: #000;
                padding: 10px;
                border-radius: 8px;
                margin-bottom: 10px;
                width: fit-content;
            }
            .message.bot {
                background-color: #10a37f;
                color: white;
                padding: 10px;
                border-radius: 8px;
                margin-bottom: 10px;
                width: fit-content;
                margin-left: auto;
            }
            .chat-container {
                border: 1px solid #ccc;
                padding: 20px;
                max-height: 400px;
                overflow-y: scroll;
                border-radius: 10px;
                background: #fdfdfd;
            }
        </style>
        """, unsafe_allow_html=True)
    else:
        st.warning("📂 Please upload some supported files to start.")

# -------------------- Page Router --------------------
if st.session_state.page == "login" or not st.session_state.is_authenticated:
    login_page()
elif st.session_state.page == "chatbot" and st.session_state.is_authenticated:
    chatbot_page()

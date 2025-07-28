import streamlit as st
from docx import Document as DocxPresentation
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import FAISS 
from langchain_community.llms import ollama
from langchain.chains import retrieval_qa
from langchain_core.documents import Document
import tempfile 
import os 

st.set_page_config(page_title = "Hikvision chatbot ",layout = "wide")
st.title("Hikvision private chatbot ")

st.markdown("""
            <style>
            .message.user{background_color : #e5e5ea; color = black; padding : 10px; border-radius: 10px, margin: 10px 0; max-width : 80% ; align-self : flex-start}
            .message.bot{background_color : #10a37f;color = whit;, padding : 10x; border-radius : 10px; margin: 10px ;max-width : 80%,  align-self: flex-end}
            .chat-container{display: flex; flex-direction: column; height : 500 px; height : 500px; overflow-y: auto, border : 1px solid #ccc ; padding : 20px; border_radius: 10px; background : #f7f7f8;}
            </style>
""", unsafe_allow_html=True)
if "history" not in st.session_state:
    st.session_state.history = []

def load_pdf(file_path):
    return PyPDFLoader(file_path).load() 
def load_docx(file_path):
    doc = DocxPresentation(file_path)
    text = "\n".join([p.text for p in doc.paragraphs])
    return [Document(page_content=text)]

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text += shape.text + "\n"
    return [Document(page_content=text)]

uploaded_files = st.file_uploader(
    "Upload PDF,PPTX or DOCX files "
)


